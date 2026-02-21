from __future__ import annotations

from pathlib import Path

from .base import ConversionResult


class PptxFallbackConverter:
    def can_convert(self, ext: str) -> bool:
        try:
            import pptx  # noqa: F401

            return ext.lower() == ".pptx"
        except ImportError:
            return False

    def convert(self, path: Path, encoding: str = "utf-8") -> ConversionResult:
        from pptx import Presentation

        prs = Presentation(str(path))
        slides: list[str] = []
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    texts.append(shape.text)
            slides.append(f"### Slide {i + 1}\n" + "\n".join(texts))
        return ConversionResult(
            content="\n\n".join(slides),
            converter_name="python-pptx",
            original_mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
